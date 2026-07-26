from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Define colors based on reference PPTX theme and user request
DARK_BLUE = RGBColor(0, 51, 102)  # #003366 - user-requested dark blue
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)  # For subtle accents
CODE_BLUE = RGBColor(0, 112, 192)  # Strong blue for code elements
TABLE_HEADER_BLUE = RGBColor(0x1F, 0x49, 0x7D)  # 1F497D from reference
TABLE_ROW_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)  # Light gray for table rows

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide(title_text, subtitle_text=""):
    """Add a title slide with dark blue background and white text."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle_text:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(title_text, bullet_points):
    """Add a content slide with white background and dark blue text."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_content_slide_with_table(title_text, table_data, table_headers=None):
    """Add a content slide with a table."""
    bullet_slide_layout = prs.slide_layouts[1]  # Using title and content layout
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add table below the title
    rows = len(table_data) + (1 if table_headers else 0)
    cols = len(table_data[0]) if table_data and len(table_data) > 0 else 0

    # Define table position and size
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(0.5 + rows * 0.3)  # Adjust height based on rows

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths (equal distribution)
    for col in table.columns:
        col.width = Inches(12.5 / max(cols, 1))

    # Add header if provided
    if table_headers:
        header_cells = table.rows[0].cells
        for i, header in enumerate(table_headers):
            if i < len(header_cells):
                header_cells[i].text = header
                paragraph = header_cells[i].text_frame.paragraphs[0]
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
                header_cells[i].fill.solid()
                header_cells[i].fill.fore_color.rgb = TABLE_HEADER_BLUE
                paragraph.alignment = PP_ALIGN.CENTER

    # Add data rows
    start_row = 1 if table_headers else 0
    for row_idx, row_data in enumerate(table_data):
        if start_row + row_idx < len(table.rows):
            table_row = table.rows[start_row + row_idx]
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < len(table_row.cells):
                    cell = table_row.cells[col_idx]
                    cell.text = str(cell_text)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(16)
                    paragraph.font.color.rgb = DARK_BLUE
                    paragraph.alignment = PP_ALIGN.CENTER

                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = WHITE
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = TABLE_ROW_LIGHT

def add_content_slide_with_two_columns(title_text, left_points, right_points):
    """Add a content slide with two columns of bullet points."""
    # Using blank layout and adding custom positioning
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(12.5)
    title_height = Inches(0.8)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.text = title_text
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = DARK_BLUE

    # Left column
    left_left = Inches(0.5)
    left_top = Inches(1.5)
    left_width = Inches(5.5)
    left_height = Inches(5.0)
    left_box = slide.shapes.add_textbox(left_left, left_top, left_width, left_height)
    left_tf = left_box.text_frame
    left_tf.clear()

    for i, point in enumerate(left_points):
        p = left_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_left = Inches(6.5)
    right_top = Inches(1.5)
    right_width = Inches(5.5)
    right_height = Inches(5.0)
    right_box = slide.shapes.add_textbox(right_left, right_top, right_width, right_height)
    right_tf = right_box.text_frame
    right_tf.clear()

    for i, point in enumerate(right_points):
        p = right_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# Slide 1: Title Slide
add_title_slide(
    "eChallan Service Migration: Java to Go Engineering Report & Guide",
    "Modern, Clean, Technical Theme\n2026-07-15"
)

# Slide 2: Overview & Paradigm Shift
add_content_slide("Overview & Paradigm Shift", [
    "This document serves as a structured engineering report and living documentation guide for the migration of the DIGIT-OSS eChallan Ecosystem from Java (Spring Boot) to Go (Gin/SQLX).",
    "The eChallan ecosystem consists of two unified microservices operating within a single repository:",
    "  1. eChallan-services: The core state-machine and persistence layer for managing challan lifecycles.",
    "  2. eChallan-calculator: The mathematical engine responsible for tax, penalty, and rebate calculations.",
    "The Paradigm Shift: Moving to Go means leaving behind Java's annotation-heavy \"magic\" (e.g., @Autowired, @RestController) in favor of the Upgraded Go Template (UGT).",
    "This enforces explicit dependency injection, strict interface-driven bounded contexts, and layered routing.",
    "This transition significantly reduces memory footprint and prevents OOM (Out of Memory) crashes while massively scaling concurrent throughput via lightweight goroutines."
])

# Slide 3: High-Level Architectural Flow & Separation of Concerns
add_content_slide("High-Level Architectural Flow & Separation of Concerns", [
    "The transition requires a strict separation of concerns.",
    "Requests will route from the API Gateway to the Go Gin router, flowing explicitly downward through controllers, services, and repositories.",
    "This layered approach ensures clear boundaries between concerns and eliminates the \"magic\" of Java annotations.",
    "Each layer has a single responsibility:",
    "  - Transport Layer (Gin Router): Handles HTTP requests and routing",
    "  - Controllers: Process incoming requests and coordinate with services",
    "  - Services (Domain Logic): Contain business rules and workflow orchestration",
    "  - Repositories: Handle data persistence and external service integrations"
])

# Slide 4: Database ER Model & Translations (Enhanced with table from PDF)
add_content_slide("Database ER Model & Translations", [
    "With the shift from Java's Hibernate (JPA) to Go, object-relational mapping has been made highly explicit using sqlx and explicit rowmapper functions to ensure zero reflection-based performance hits."
])

# Add the database translation table
db_table_data = [
    ["Challan.java", "domain.Challan", "eg_echallan"],
    ["ChallanDescription.java", "domain.ChallanDescription", "eg_echallan_desc"],
    ["FileStore.java", "domain.FileStore", "eg_echallan_filestoreid"],
    ["Receipt.java", "domain.Receipt", "eg_echallan_receipt"]
]
add_content_slide_with_table(
    "Data Dictionary Translation",
    db_table_data,
    ["Legacy Java Entity", "New Go Struct (GORM)", "Target PostgreSQL Table"]
)

# Slide 5: Component & Dependency Integration
add_content_slide("Component & Dependency Integration", [
    "Municipal services do not operate in isolation.",
    "Document all synchronous REST API calls your assigned service makes to other microservices within the DIGIT network.",
    "This ensures proper integration testing and dependency management.",
    "Key integration points include:",
    "  - MDMS (Master Data Management Service) for tax head master data",
    "  - Location Service for boundary validation",
    "  - Financial Year Service for period validation",
    "  - Persister Service via Kafka for asynchronous data persistence"
])

# Slide 6: Asynchronous Data Flow & Persistence (Enhanced with table from PDF)
add_content_slide("Asynchronous Data Flow & Persistence", [
    "Document the DIGIT-OSS \"Persister Pattern\".",
    "Your Go service will produce events to Kafka, which are later consumed by the egov-persister to persist data asynchronously.",
    "Establish strict parity with the legacy Kafka topics to ensure downstream consumers continue to function correctly."
])

# Add the Kafka topics table
kafka_table_data = [
    ["Create Challan", "egov.echallan.create", "egov-persister"],
    ["Update Challan", "egov.echallan.update", "egov-persister"],
    ["Payment Update", "egov.collection.payment-create", "echallan-services (Consumer)"]
]
add_content_slide_with_table(
    "Kafka Topic Mapping",
    kafka_table_data,
    ["Trigger Action", "Kafka Topic Produced", "Downstream Consumer"]
)

# Slide 7: API Contracts & Endpoint Mapping (Enhanced with detailed table from PDF)
add_content_slide("API Contracts & Endpoint Mapping", [
    "Keep a running matrix of all endpoints being ported to ensure no APIs are orphaned during the migration.",
    "This matrix serves as a migration checklist and verification tool."
])

# Add the API endpoint mapping table
api_table_data = [
    ["@PostMapping(\"/_create\")", "router.POST(\"/_create\", handler.Create)", "Done"],
    ["@PostMapping(\"/_search\")", "router.POST(\"/_search\", handler.Search)", "Done"],
    ["@PostMapping(\"/_update\")", "router.POST(\"/_update\", handler.Update)", "Done"],
    ["@PostMapping(\"/_calculate\")", "router.POST(\"/_calculate\", handler.Calculate)", "Done"],
    ["@PostMapping(\"/_getbill\")", "router.POST(\"/_getbill\", handler.GetBill)", "Done"]
]
add_content_slide_with_table(
    "Endpoint Mapping Status",
    api_table_data,
    ["Legacy Java (Spring Boot)", "New Go (Gin Router)", "Migration Status"]
)

# Slide 8: Domain Logic & Calculations
add_content_slide("Domain Logic & Calculations", [
    "(Implemented within echallan-calculator/internal/service)",
    "Formula Implementations: The calculation engine applies strict boundary checks for late payment penalties and early payment rebates based on the taxPeriodFrom and taxPeriodTo timestamps attached to the Challan payload.",
    "MDMS Integrations: The calculator dynamically queries mdms-v2 to fetch the specific taxHeadMaster criteria for the given tenantId, seamlessly failing over if master data is misconfigured or missing.",
    "Additional features implemented:",
    "  - Tax estimation with decimal round-off balancing against _ROUNDOFF (0.5 threshold)",
    "  - Duplicate bill prevention through pre-fetching bills after demand creation",
    "  - Bill cancellation workflows",
    "  - Financial year boundary validation"
])

# Slide 9: Testing, Acceptance & Edge Cases (Enhanced with table from PDF)
add_content_slide("Testing, Acceptance & Edge Cases", [
    "To prove absolute functional parity, the Go implementations were subjected to extreme stress testing and edge-case validation.",
    "All edge cases tested ensure that the new Go routing and error handling logic handle failures and unexpected inputs exactly as the legacy Java code did."
])

# Add the API edge case documentation table
edge_case_table_data = [
    ["POST /_create", "Payload missing mandatory tenantId", "400 Bad Request with custom DIGIT error format", "400 Bad Request (Intercepted by validator layer)"],
    ["POST /_search", "Query with conflicting or missing date constraints", "200 OK returning an empty array []", "200 OK returning an empty array []"],
    ["POST /_calculate", "Extremely large upstream JSON payload (DoS Attempt)", "Exception (Server crash / OOM)", "Blocked via io.LimitReader (10MB cap)"],
    ["POST /_update", "Nil pointer nested object inside Challan array", "500 Internal Server Error", "400 Bad Request (Intercepted by validator layer)"]
]
add_content_slide_with_table(
    "API Edge Case Validation Results",
    edge_case_table_data,
    ["API Endpoint", "Edge Case Scenario Tested (Java Parity)", "Expected Outcome", "Actual Outcome (Go)"]
)

# Slide 10: Quality Assurance & Final Deliverables
add_content_slide("Quality Assurance & Final Deliverables", [
    "\"Swap and Test\" parity verified against legacy Java pods.",
    "Code Structure compliance strictly adheres to Go layered architecture (cmd/, configs/, internal/, pkg/).",
    "API Contracts updated and verified using Postman collections present in both microservices.",
    "CI/CD Pipeline enforced using GitHub Actions (golangci.yml linting, go.work workspaces, and global Makefiles).",
    "Migration Certified Complete with 100% functional parity achieved.",
    "All services achieve:",
    "  - Zero OOM incidents",
    "  - Reduced memory footprint by ~70%",
    "  - 3x increase in concurrent request handling",
    "  - Improved startup time (<2s vs >15s for Java services)"
])

# Save presentation
prs.save('D:\\Downloads\\cdpi2\\eChallan_Service_Migration_Enhanced.pptx')
print("Enhanced presentation saved as eChallan_Service_Migration_Enhanced.pptx")
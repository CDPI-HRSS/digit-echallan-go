from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_ARROWHEAD
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

# Define colors based on reference PPTX theme and user request
DARK_BLUE = RGBColor(0, 51, 102)  # #003366 - user-requested dark blue
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)  # For subtle accents
CODE_BLUE = RGBColor(0, 112, 192)  # Strong blue for code elements
TABLE_HEADER_BLUE = RGBColor(0x1F, 0x49, 0x7D)  # 1F497D from reference
TABLE_ROW_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)  # Light gray for table rows
ACCENT_BLUE = RGBColor(0x4F, 0x81, 0xBD)  # 4F81BD - accent1
ACCENT_GREEN = RGBColor(0x9B, 0xBB, 0x59)  # 9BBB59 - accent3
ACCENT_ORANGE = RGBColor(0xF7, 0x96, 0x46)  # F79646 - accent6

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide(title_text, subtitle_text=""):
    """Add a title slide with dark blue background and white text."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle_text:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(title_text, bullet_points):
    """Add a content slide with white background and dark blue text."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_content_slide_with_table(title_text, table_data, table_headers=None):
    """Add a content slide with a table."""
    bullet_slide_layout = prs.slide_layouts[1]  # Using title and content layout
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add table below the title
    rows = len(table_data) + (1 if table_headers else 0)
    cols = len(table_data[0]) if table_data and len(table_data) > 0 else 0

    # Define table position and size - make it smaller to leave room for potential diagrams
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.0)  # Slightly narrower
    height = Inches(0.5 + rows * 0.3)  # Adjust height based on rows

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths (equal distribution)
    for col in table.columns:
        col.width = Inches(12.0 / max(cols, 1))

    # Add header if provided
    if table_headers:
        header_cells = table.rows[0].cells
        for i, header in enumerate(table_headers):
            if i < len(header_cells):
                header_cells[i].text = header
                paragraph = header_cells[i].text_frame.paragraphs[0]
                paragraph.font.size = Pt(16)  # Slightly smaller
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
                header_cells[i].fill.solid()
                header_cells[i].fill.fore_color.rgb = TABLE_HEADER_BLUE
                paragraph.alignment = PP_ALIGN.CENTER

    # Add data rows
    start_row = 1 if table_headers else 0
    for row_idx, row_data in enumerate(table_data):
        if start_row + row_idx < len(table.rows):
            table_row = table.rows[start_row + row_idx]
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < len(table_row.cells):
                    cell = table_row.cells[col_idx]
                    cell.text = str(cell_text)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(14)  # Smaller for tables
                    paragraph.font.color.rgb = DARK_BLUE
                    paragraph.alignment = PP_ALIGN.CENTER

                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = WHITE
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = TABLE_ROW_LIGHT

def add_flow_diagram_slide(title_text, flow_steps, flow_labels=None, is_horizontal=True):
    """Add a slide with a simple flow diagram using shapes."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(12.5)
    title_height = Inches(0.8)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.text = title_text
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = DARK_BLUE

    if is_horizontal:
        # Horizontal flow
        step_width = Inches(2.0)
        step_height = Inches(1.0)
        start_x = Inches(1.0)
        y_pos = Inches(2.5)

        for i, step in enumerate(flow_steps):
            x_pos = start_x + (i * (step_width + Inches(0.5)))

            # Draw step box
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                x_pos, y_pos, step_width, step_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT_BLUE
            shape.line.color.rgb = DARK_BLUE

            # Add text to shape
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.text = step
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                p = shape.text_frame.add_paragraph()
                p.text = step
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

            # Add label if provided
            if flow_labels and i < len(flow_labels):
                label_left = x_pos
                label_top = y_pos + step_height + Inches(0.1)
                label_width = step_width
                label_height = Inches(0.3)
                label_box = slide.shapes.add_textbox(label_left, label_top, label_width, label_height)
                label_tf = label_box.text_frame
                label_tf.text = flow_labels[i]
                label_tf.paragraphs[0].font.size = Pt(12)
                label_tf.paragraphs[0].font.color.rgb = DARK_BLUE
                label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Add arrow to next step (except for last step)
            if i < len(flow_steps) - 1:
                arrow_start_x = x_pos + step_width
                arrow_end_x = arrow_start_x + Inches(0.5)
                arrow_y = y_pos + step_height/2

                arrow = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    arrow_start_x, arrow_y, arrow_end_x, arrow_y
                )
                arrow.line.width = Pt(2.5)
                arrow.line.color.rgb = DARK_BLUE
                # Arrow head
                arrow.line.end_arrowhead = MSO_ARROWHEAD.TRIANGLE
    else:
        # Vertical flow
        step_width = Inches(4.0)
        step_height = Inches(0.8)
        start_y = Inches(2.0)
        x_pos = Inches(4.5)

        for i, step in enumerate(flow_steps):
            y_pos = start_y + (i * (step_height + Inches(0.3)))

            # Draw step box
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                x_pos, y_pos, step_width, step_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT_BLUE
            shape.line.color.rgb = DARK_BLUE

            # Add text to shape
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.text = step
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                p = shape.text_frame.add_paragraph()
                p.text = step
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

            # Add label if provided
            if flow_labels and i < len(flow_labels):
                label_left = x_pos + step_width + Inches(0.2)
                label_top = y_pos
                label_width = Inches(2.0)
                label_height = step_height
                label_box = slide.shapes.add_textbox(label_left, label_top, label_width, label_height)
                label_tf = label_box.text_frame
                label_tf.text = flow_labels[i]
                label_tf.paragraphs[0].font.size = Pt(12)
                label_tf.paragraphs[0].font.color.rgb = DARK_BLUE
                label_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Add arrow to next step (except for last step)
            if i < len(flow_steps) - 1:
                arrow_start_y = y_pos + step_height
                arrow_end_y = arrow_start_y + Inches(0.3)
                arrow_x = x_pos + step_width/2

                arrow = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    arrow_x, arrow_start_y, arrow_x, arrow_end_y
                )
                arrow.line.width = Pt(2.5)
                arrow.line.color.rgb = DARK_BLUE
                # Arrow head
                arrow.line.end_arrowhead = MSO_ARROWHEAD.TRIANGLE

def add_themed_content_slide(title_text, content_elements):
    """Add a slide with themed content including icons or visual indicators."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, element in enumerate(content_elements):
        # Add a small colored square as a visual indicator
        p = tf.add_paragraph()
        p.text = "  " + element  # Indent with spaces
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# Slide 1: Title Slide
add_title_slide(
    "eChallan Service Migration: Java to Go Engineering Report & Guide",
    "Modern, Clean, Technical Theme\n2026-07-15"
)

# Slide 2: Overview & Paradigm Shift with simple diagram
overview_steps = [
    "Legacy Java Spring Boot",
    "Migration Process",
    "Modern Go Gin/SQLX Stack"
]
overview_labels = [
    "Annotation-heavy 'magic'",
    "Explicit, Layered UGT Approach",
    "Reduced OOM, High Throughput"
]
add_flow_diagram_slide("Overview & Paradigm Shift", overview_steps, overview_labels, is_horizontal=True)

# Slide 3: High-Level Architectural Flow
flow_steps = [
    "API Gateway",
    "Go Gin Router\n(Transport Layer)",
    "Controllers",
    "Services\n(Domain Logic)",
    "Repositories\n(PostgreSQL & Kafka)"
]
add_flow_diagram_slide("High-Level Architectural Flow & Separation of Concerns", flow_steps, is_horizontal=False)

# Slide 4: Database ER Model & Translations
add_content_slide("Database ER Model & Translations", [
    "We transitioned from Java's Hibernate (JPA) to highly explicit Go 'sqlx' mapping and explicit rowmapper functions.",
    "This ensures zero reflection-based performance hits and gives us full control over database operations."
])

# Add the database translation table
db_table_data = [
    ["Challan.java", "domain.Challan", "eg_echallan"],
    ["ChallanDescription.java", "domain.ChallanDescription", "eg_echallan_desc"],
    ["FileStore.java", "domain.FileStore", "eg_echallan_filestoreid"],
    ["Receipt.java", "domain.Receipt", "eg_echallan_receipt"]
]
add_content_slide_with_table(
    "Data Dictionary Translation",
    db_table_data,
    ["Legacy Java Entity", "New Go Struct (GORM)", "Target PostgreSQL Table"]
)

# Slide 5: Component Dependencies with visual indicators
deps_content = [
    "MDMS (Master Data Management Service) - Tax Head Master Data",
    "Location Service - Boundary Validation",
    "Financial Year Service - Period Validation",
    "Kafka Persister Service - Asynchronous Data Persistence",
    "Internal Service-to-Service Communication via Defined APIs"
]
add_themed_content_slide("Component & Dependency Integration", deps_content)

# Slide 6: Asynchronous Data Flow & Persistence
add_content_slide("Asynchronous Data Flow & Persistence", [
    "Implementation of the DIGIT-OSS \"Persister Pattern\" for decoupled persistence.",
    "Go services produce events to Kafka for asynchronous processing by egov-persister."
])

# Add the Kafka topics table
kafka_table_data = [
    ["Create Challan", "egov.echallan.create", "egov-persister"],
    ["Update Challan", "egov.echallan.update", "egov-persister"],
    ["Payment Update", "egov.collection.payment-create", "echallan-services (Consumer)"]
]
add_content_slide_with_table(
    "Kafka Topic Mapping",
    kafka_table_data,
    ["Trigger Action", "Kafka Topic Produced", "Downstream Consumer"]
)

# Slide 7: API Contracts & Endpoint Mapping
api_table_data = [
    ["@PostMapping(\"/_create\")", "router.POST(\"/_create\", handler.Create)", "Done"],
    ["@PostMapping(\"/_search\")", "router.POST(\"/_search\", handler.Search)", "Done"],
    ["@PostMapping(\"/_update\")", "router.POST(\"/_update\", handler.Update)", "Done"],
    ["@PostMapping(\"/_calculate\")", "router.POST(\"/_calculate\", handler.Calculate)", "Done"],
    ["@PostMapping(\"/_getbill\")", "router.POST(\"/_getbill\", handler.GetBill)", "Done"]
]
add_content_slide_with_table(
    "Endpoint Mapping Status",
    api_table_data,
    ["Legacy Java (Spring Boot)", "New Go (Gin Router)", "Migration Status"]
)

# Slide 8: Domain Logic & Calculations
calc_content = [
    "✓ Strict boundary checks for late payment penalties & early payment rebates",
    "✓ Based on taxPeriodFrom and taxPeriodTo timestamps in payload",
    "✓ Dynamic MDMS v2 queries for taxHeadMaster criteria",
    "✓ Clean failover if master data is misconfigured or missing",
    "✓ Tax estimation with decimal round-off balancing (_ROUNDOFF = 0.5 threshold)",
    "✓ Duplicate bill prevention through pre-fetching after demand creation",
    "✓ Complete bill cancellation workflows",
    "✓ Financial year boundary validation"
]
add_themed_content_slide("Domain Logic & Calculations", calc_content)

# Slide 9: Testing, Acceptance & Edge Cases
edge_case_table_data = [
    ["POST /_create", "Missing tenantId", "400 Bad Request", "400 Bad Request (Validator)"],
    ["POST /_search", "Conflicting/missing dates", "200 OK + empty array []", "200 OK + empty array []"],
    ["POST /_calculate", "Large JSON payload (DoS)", "Exception/OOM crash", "Blocked by io.LimitReader (10MB)"],
    ["POST /_update", "Nil pointer in array", "500 Internal Server Error", "400 Bad Request (Validator)"]
]
add_content_slide_with_table(
    "API Edge Case Validation Results",
    edge_case_table_data,
    ["API Endpoint", "Edge Case Scenario", "Expected Outcome (Java)", "Actual Outcome (Go)"]
)

# Slide 10: Quality Assurance & Final Deliverables with metrics diagram
qa_content = [
    "✓ \"Swap and Test\" parity verified against legacy Java pods",
    "✓ Code Structure: Go layered architecture (cmd/, configs/, internal/, pkg/)",
    "✓ API Contracts: Updated & verified with Postman collections",
    "✓ CI/CD: GitHub Actions (golangci.yml, go.work, Makefiles)",
    "✓ Migration Certified Complete"
]
add_themed_content_slide("Quality Assurance & Final Deliverables", qa_content)

# Slide 11: Performance Improvements (using simple visual)
perf_title = "Performance Improvements Achieved"
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
title_tf = title_box.text_frame
title_tf.text = perf_title
title_tf.paragraphs[0].font.size = Pt(32)
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].font.color.rgb = DARK_BLUE

# Performance metrics with visual indicators
metrics = [
    ("Memory Footprint", "~70% Reduction", ACCENT_GREEN),
    ("Concurrent Requests", "3x Increase", ACCENT_BLUE),
    ("Startup Time", "<2s vs >15s", ACCENT_ORANGE),
    ("OOM Incidents", "Zero Occurrences", ACCENT_GREEN)
]

# Arrange metrics in 2x2 grid
positions = [
    (Inches(1.5), Inches(2.0)),
    (Inches(7.5), Inches(2.0)),
    (Inches(1.5), Inches(4.0)),
    (Inches(7.5), Inches(4.0))
]

for i, (label, value, color) in enumerate(metrics):
    left, top = positions[i]

    # Draw metric box
    box_width = Inches(5.0)
    box_height = Inches(1.5)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, box_width, box_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(1.5)

    # Add label
    label_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), box_width - Inches(0.4), Inches(0.4))
    label_tf = label_box.text_frame
    label_tf.text = label
    label_tf.paragraphs[0].font.size = Pt(16)
    label_tf.paragraphs[0].font.bold = True
    label_tf.paragraphs[0].font.color.rgb = DARK_BLUE

    # Add value
    value_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), box_width - Inches(0.4), Inches(0.4))
    value_tf = value_box.text_frame
    value_tf.text = value
    value_tf.paragraphs[0].font.size = Pt(20)
    value_tf.paragraphs[0].font.bold = True
    value_tf.paragraphs[0].font.color.rgb = color

# Slide 12: Migration Summary
summary_content = [
    "🎯 Successfully migrated from Java Spring Boot to Go Gin/SQLX",
    "🔧 Applied Upgraded Go Template (UGT) for explicit, layered architecture",
    "⚡ Achieved significant performance improvements and reliability",
    "🛡️ Maintained 100% functional parity with comprehensive testing",
    "📦 Established robust CI/CD pipeline and deployment practices"
]
add_themed_content_slide("Migration Summary & Success Metrics", summary_content)

# Save presentation
prs.save('D:\\Downloads\\cdpi2\\eChallan_Service_Migration_With_Diagrams.pptx')
print("Presentation with diagrams saved as eChallan_Service_Migration_With_Diagrams.pptx")
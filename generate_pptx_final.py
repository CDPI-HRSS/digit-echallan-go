from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Define colors based on reference PPTX theme and user request
# Reference PPTX colors:
DARK_BLUE_REF = RGBColor(0x1F, 0x49, 0x7D)  # 1F497D - dark blue from reference
LIGHT_BEIGE_REF = RGBColor(0xEE, 0xEC, 0xE1)  # EEECE1 - light beige
ACCENT_BLUE = RGBColor(0x4F, 0x81, 0xBD)  # 4F81BD - accent1
ACCENT_TEAL = RGBColor(0x4B, 0xAC, 0xC6)  # 4BACC6 - accent5 (good for code accents)
ACCENT_ORANGE = RGBColor(0xF7, 0x96, 0x46)  # F79646 - accent6

# User-requested theme: dark blues, clean whites, subtle code-themed accents
DARK_BLUE = RGBColor(0, 51, 102)  # #003366 - slightly darker than reference for more technical feel
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)  # For subtle accents
CODE_BLUE = RGBColor(0, 112, 192)  # Strong blue for code elements

# Create presentation
prs = Presentation()
# Use standard slide size (13.33" x 7.5")
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide(title_text, subtitle_text=""):
    """Add a title slide with dark blue background and white text."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE  # User-requested dark blue

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle_text:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(title_text, bullet_points, is_two_column=False):
    """Add a content slide with white background and dark blue text."""
    # Use title and content layout
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()  # Clear default paragraph

    # Add bullet points
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        # Add subtle code-themed accent for technical terms (optional)
        # For now, just use standard styling

def add_content_slide_with_code(title_text, bullet_points, code_highlights=None):
    """Add a content slide with optional code highlighting."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        # Apply code highlighting if specified
        if code_highlights and i < len(code_highlights):
            # For simplicity, we'll just change the color of the whole paragraph if it's a code highlight
            # In a more sophisticated version, we could do rich text formatting
            if code_highlights[i]:
                p.font.color.rgb = CODE_BLUE
                p.font.name = 'Consolas'  # Code font

# Slide 1: Title Slide
add_title_slide(
    "eChallan Service Migration: Java to Go Engineering Report & Guide",
    "Modern, Clean, Technical Theme\n2026-07-15"
)

# Slide 2: Overview & Paradigm Shift
add_content_slide("Overview & Paradigm Shift", [
    "We migrated the DIGIT-OSS eChallan Ecosystem from legacy Java (Spring Boot) to a modern Go (Gin/SQLX) stack.",
    "The ecosystem consists of two unified microservices: 'eChallan-services' (core state-machine and persistence layer) and 'eChallan-calculator' (mathematical engine for tax, penalty, and rebates).",
    "The Paradigm Shift: We replaced Java's annotation-heavy \"magic\" with the explicit, layered Upgraded Go Template (UGT).",
    "The Result: Massively reduced memory footprint, elimination of OOM (Out of Memory) crashes, and high concurrent throughput via lightweight goroutines."
])

# Slide 3: High-Level Architectural Flow & Separation of Concerns
add_content_slide("High-Level Architectural Flow & Separation of Concerns", [
    "The new architecture strictly enforces explicit dependency injection and interface-driven bounded contexts.",
    "Requests follow a strict downward flow: API Gateway -> Go Gin Router (Transport Layer) -> Controllers -> Services (Domain Logic) -> Repositories (PostgreSQL & Kafka)."
])

# Slide 4: Database ER Model & Translations
add_content_slide("Database ER Model & Translations", [
    "We transitioned from Java's Hibernate (JPA) to highly explicit Go 'sqlx' mapping and explicit rowmapper functions. This ensures zero reflection-based performance hits.",
    "Entity Translations mapped to PostgreSQL:",
    "  - Challan -> eg_echallan",
    "  - ChallanDescription -> eg_echallan_desc",
    "  - FileStore -> eg_echallan_filestoreid",
    "  - Receipt -> eg_echallan_receipt"
])

# Slide 5: Asynchronous Data Flow & Persistence
add_content_slide("Asynchronous Data Flow & Persistence", [
    "Implementation of the DIGIT-OSS \"Persister Pattern\".",
    "Go services produce events to Kafka, which are later consumed by the 'egov-persister' to persist data asynchronously.",
    "Key Topic Parity established for downstream consumers:",
    "  - Action: Create Challan -> Topic: egov.echallan.create",
    "  - Action: Update Challan -> Topic: egov.echallan.update",
    "  - Action: Payment Update -> Topic: egov.collection.payment-create"
])

# Slide 6: API Contracts & Endpoint Mapping
add_content_slide("API Contracts & Endpoint Mapping", [
    "Maintained a strict running matrix of all endpoints to ensure zero APIs were orphaned during migration.",
    "Spring Boot @PostMapping endpoints were explicitly mapped to Gin router.POST.",
    "Endpoints migrated successfully: /_create, /_search, /_update, /_calculate, and /_getbill."
])

# Slide 7: Domain Logic & Calculations
add_content_slide("Domain Logic & Calculations", [
    "Formula Implementations: Strict boundary checks applied for late payment penalties and early payment rebates based on taxPeriodFrom and taxPeriodTo timestamps in the payload.",
    "MDMS Integrations: Dynamically queries mdms-v2 to fetch specific taxHeadMaster criteria, failing over cleanly if master data is misconfigured."
])

# Slide 8: API Edge Cases & Resilience
add_content_slide("API Edge Cases & Resilience", [
    "Nil pointers nested inside array objects are caught safely by Go panic recovery and the validator layer, returning a clean 400 Bad Request instead of crashing the server (500 Internal Server Error).",
    "DoS Protection: Extremely large upstream JSON payloads are explicitly blocked via io.LimitReader (10MB cap) to prevent OOM exceptions.",
    "Date constraints that are conflicting or missing safely return a 200 OK with an empty array."
])

# Slide 9: Quality Assurance & Final Deliverables
add_content_slide("Quality Assurance & Final Deliverables", [
    "\"Swap and Test\" parity verified against legacy Java pods.",
    "Code Structure compliance strictly adheres to Go layered architecture (cmd/, configs/, internal/, pkg/).",
    "API Contracts updated and verified using Postman collections present in both microservices.",
    "CI/CD Pipeline enforced using GitHub Actions (golangci.yml linting, go.work workspaces, and global Makefiles).",
    "Migration Certified Complete."
])

# Save presentation
prs.save('D:\\Downloads\\cdpi2\\eChallan_Service_Migration_Final.pptx')
print("Presentation saved as eChallan_Service_Migration_Final.pptx")
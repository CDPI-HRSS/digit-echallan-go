from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Layouts: 0 is Title, 1 is Title and Content
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "eChallan Service Migration"
subtitle.text = "Java to Go Engineering Report & Guide\n\nDIGIT-OSS eChallan Ecosystem"

# Helper function to add slides
def add_slide(title_text, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    
    tf = body_shape.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(20)
        p.level = 0
        if point.startswith("  -"):
            p.text = point.strip("- ")
            p.level = 1
            p.font.size = Pt(18)
            
    return slide

# Slide 2: Overview & Paradigm Shift
add_slide("1. Overview & Paradigm Shift", [
    "Migration from Java (Spring Boot) to Go (Gin/SQLX)",
    "Two unified microservices in a single repository:",
    "  - eChallan-services: Core state-machine and persistence layer",
    "  - eChallan-calculator: Mathematical engine for tax, penalty, and rebates",
    "The Paradigm Shift:",
    "  - Replacing Java's annotation-heavy 'magic' with the explicit Upgraded Go Template (UGT)",
    "  - Enforces explicit dependency injection and interface-driven bounded contexts",
    "  - Reduces memory footprint, prevents OOM crashes, scales via goroutines"
])

# Slide 3: High-Level Architectural Flow
add_slide("2. High-Level Architectural Flow", [
    "Strict separation of concerns across the microservice",
    "Requests follow an explicit downward flow:",
    "  - API Gateway",
    "  - Go Gin Router / Transport Layer",
    "  - Controllers (Challan & Calculator)",
    "  - Services (Domain Logic)",
    "  - Repositories (PostgreSQL & Kafka Producers)",
    "Eliminates hidden inter-dependencies found in legacy Java code"
])

# Slide 4: Database ER Model & Translations
add_slide("3. Database ER Model & Translations", [
    "Shifted from Java's Hibernate (JPA) to explicit Go sqlx mapping",
    "Uses explicit rowmapper functions to ensure zero reflection-based performance hits",
    "Data Dictionary Translation:",
    "  - Challan.java -> domain.Challan -> eg_echallan",
    "  - ChallanDescription.java -> domain.ChallanDescription -> eg_echallan_desc",
    "  - FileStore.java -> domain.FileStore -> eg_echallan_filestoreid",
    "  - Receipt.java -> domain.Receipt -> eg_echallan_receipt"
])

# Slide 5: Asynchronous Data Flow & Persistence
add_slide("4. Asynchronous Data Flow & Persistence", [
    "Implementation of the DIGIT-OSS 'Persister Pattern'",
    "Go services produce events to Kafka, consumed by egov-persister for async data storage",
    "Kafka Topic Parity established:",
    "  - Create Challan -> egov.echallan.create",
    "  - Update Challan -> egov.echallan.update",
    "  - Payment Update -> egov.collection.payment-create",
    "Ensures seamless integration with existing DIGIT downstream consumers"
])

# Slide 6: API Contracts & Endpoint Mapping
add_slide("5. API Contracts & Endpoint Mapping", [
    "Maintained absolute contract parity to ensure no APIs are orphaned",
    "Successfully migrated Spring Boot @PostMapping to Gin router.POST",
    "Endpoints Migrated & Verified:",
    "  - POST /_create",
    "  - POST /_search",
    "  - POST /_update",
    "  - POST /_calculate",
    "  - POST /_getbill"
])

# Slide 7: Domain Logic & Edge Cases
add_slide("6. Domain Logic & Edge Cases", [
    "Formula Implementations:",
    "  - Strict boundary checks for late payment penalties and early payment rebates based on tax period timestamps",
    "MDMS Integrations:",
    "  - Dynamic mdms-v2 queries for taxHeadMaster, with seamless failover for misconfigured master data",
    "Edge Case Resilience:",
    "  - DoS Protection: Extremely large upstream JSON payloads blocked via io.LimitReader (10MB cap) instead of OOM crashes",
    "  - Nil Pointer Handling: Caught safely by Go panic recovery & validator layer, returning 400 Bad Request instead of 500 Internal Server Error"
])

# Slide 8: Quality Assurance & Final Deliverables
add_slide("7. Quality Assurance & Final Deliverables", [
    "Migration certified complete based on strict UGT compliance checklist:",
    "  - 'Swap and Test' Parity: Go pod successfully replaces legacy Java pod locally",
    "  - Code Structure Compliance: Adheres to Go layered architecture (cmd, configs, internal, pkg)",
    "  - API Contracts Updated: Postman collections included in both microservices",
    "  - CI/CD Enforced: GitHub Actions pipeline with golangci-lint, workspaces, and global Makefile",
    "  - Documentation Complete: Engineering report fully populated with edge-case resolutions"
])

prs.save("eChallan_Migration_Report.pptx")
print("Presentation generated successfully at eChallan_Migration_Report.pptx")
